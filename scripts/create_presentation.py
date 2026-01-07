import json
import os
from pptx import Presentation
from pptx.util import Inches


def build_presentation():
    root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
    results_dir = os.path.join(root, 'results')
    metrics_path = os.path.join(results_dir, 'evaluation_metrics.json')

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Clustering Evaluation"
    try:
        subtitle = slide.placeholders[1]
        subtitle.text = "Generated from results/evaluation_metrics.json"
    except Exception:
        pass

    # Table of Contents
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Table of Contents"
    try:
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, sec in enumerate(["Evaluation Overview", "Metrics", "Visuals", "Conclusions"], start=1):
            p = body.add_paragraph()
            p.text = f"{i}. {sec}"
            p.level = 0
    except Exception:
        pass

    # Metrics slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Metrics"
    try:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        if os.path.exists(metrics_path):
            with open(metrics_path, 'r') as fh:
                metrics = json.load(fh)
            for k in ['silhouette', 'calinski_harabasz', 'davies_bouldin', 'adjusted_rand', 'adjusted_mutual_info', 'purity', 'n_samples', 'n_clusters']:
                v = metrics.get(k)
                p = tf.add_paragraph()
                p.text = f"{k}: {v}"
                p.level = 0
        else:
            tf.text = "No metrics file found."
    except Exception:
        pass

    # Visuals slides
    image_files = ['cluster_pca.png', 'cluster_quality_boxplot.png', 'cluster_quality_crosstab.png', 'silhouette_distribution.png']
    for img in image_files:
        img_path = os.path.join(results_dir, img)
        if os.path.exists(img_path):
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5])
                left = Inches(0.5)
                top = Inches(1)
                max_width = Inches(9)
                prs.slide_height = prs.slide_height
                slide.shapes.add_picture(img_path, left, top, width=max_width)
            except Exception:
                pass

    # Conclusions
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusions"
    try:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.text = "Key takeaways:"
        p = tf.add_paragraph()
        p.text = "- Silhouette ~0.21 → moderate separation"
        p = tf.add_paragraph()
        p.text = "- Low ARI / AMI / Purity → weak alignment with labels"
    except Exception:
        pass

    out_path = os.path.join(results_dir, 'presentation.pptx')
    prs.save(out_path)
    print(f"Saved presentation to {out_path}")


if __name__ == '__main__':
    build_presentation()
